import os
import sys
import tkinter as tk
from tkinter import messagebox, simpledialog, ttk
import time
import base64
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches
import shutil
from selenium import webdriver
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from selenium.webdriver.edge.service import Service
from selenium.webdriver.edge.options import Options


# pappt_debug_snippets.py
# 调试辅助函数：用于检查当前窗口/文档状态、保存调试页面、用 JS 校验/点击 XPath、穿透 shadow DOM 等。
# 直接将需要的函数复制到你的 crawl_ppt 中相应位置，按注释调用。

from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from selenium.webdriver.common.by import By
import time
import os

from selenium.common.exceptions import (
    StaleElementReferenceException,
    NoSuchElementException,
    ElementClickInterceptedException,
    TimeoutException
)
from selenium.webdriver.support import expected_conditions as EC
from selenium.webdriver.common.by import By

def locate_replay_parent(self, time_text, teacher_text, timeout=8):
    """
    重新在当前页面查找与 replay_item 对应的 <p> 元素（返回 Selenium WebElement 或 None）。
    使用 time_text 和 teacher_text 做匹配（你也可以改成更可靠的匹配规则）。
    """
    end = time.time() + timeout
    while time.time() < end:
        try:
            parents = self.driver.find_elements(By.CSS_SELECTOR, ".content-inner-one > p")
            for p in parents:
                try:
                    spans = p.find_elements(By.TAG_NAME, "span")
                    if len(spans) >= 3:
                        t = spans[1].text.strip()
                        teach = spans[2].get_attribute("title") or spans[2].text.strip()
                        if t == time_text and teach == teacher_text:
                            return p
                except Exception:
                    continue
        except Exception:
            pass
        time.sleep(0.3)
    return None

def safe_click_span_button(self, parent_p, span_index=4, timeout=6):
    """
    在 parent_p（<p>）内部安全点击第 span_index 个 span（从0开始计）。
    会重试并在必要时使用 JS click。
    返回 True/False。
    """
    end = time.time() + timeout
    while time.time() < end:
        try:
            spans = parent_p.find_elements(By.TAG_NAME, "span")
            if len(spans) > span_index:
                btn = spans[span_index]
                try:
                    # 优先正常点击
                    btn.click()
                    return True
                except (ElementClickInterceptedException, StaleElementReferenceException):
                    # 尝试 JS 点击
                    try:
                        self.driver.execute_script("arguments[0].scrollIntoView(true);arguments[0].click();", btn)
                        return True
                    except Exception:
                        pass
            else:
                # 子元素数量不够，可能 DOM 还没完全渲染，短暂等待再试
                time.sleep(0.2)
        except StaleElementReferenceException:
            # parent_p 已 stale，外面逻辑应重新 locate
            return False
        except Exception:
            time.sleep(0.2)
    return False

def wait_for_ready(driver, timeout=30):
    """等待 document.readyState == 'complete'，返回 True/False。"""
    try:
        WebDriverWait(driver, timeout).until(lambda d: d.execute_script("return document.readyState") == "complete")
        return True
    except Exception:
        return False


def save_debug(driver, prefix="debug"):
    """保存截图和 HTML，文件保存在当前工作目录。
    返回 (screenshot_path, html_path)
    """
    ts = int(time.time())
    png = f"{prefix}_{ts}.png"
    html = f"{prefix}_{ts}.html"
    try:
        driver.save_screenshot(png)
    except Exception as e:
        png = None
    try:
        with open(html, "w", encoding="utf-8") as f:
            f.write(driver.page_source)
    except Exception as e:
        html = None
    return png, html


def xpath_exists_js(driver, xpath):
    """在当前 document 上用 document.evaluate 检查 XPath 是否存在（返回 True/False）。"""
    script = (
        "var r=document.evaluate(arguments[0], document, null, XPathResult.FIRST_ORDERED_NODE_TYPE, null);"
        "return r.singleNodeValue !== null;"
    )
    try:
        return bool(driver.execute_script(script, xpath))
    except Exception:
        return False


def get_element_by_xpath_js(driver, xpath):
    """用 JS 返回 XPath 对应的 DOM 节点（可传回给 execute_script 并与 arguments[0] 一起使用）。
    注意：直接返回的对象在 Selenium 层面是 WebElement 的 JS 引用，可用于 arguments[0] 操作。
    若找不到返回 None。
    """
    script = (
        "var r=document.evaluate(arguments[0], document, null, XPathResult.FIRST_ORDERED_NODE_TYPE, null);"
        "return r.singleNodeValue;"
    )
    try:
        return driver.execute_script(script, xpath)
    except Exception:
        return None


def click_via_js(driver, element_or_xpath):
    """如果传入 WebElement（Selenium 返回的元素引用）或 XPath 字符串，使用 JS 点击。
    返回 True/False。
    """
    try:
        if isinstance(element_or_xpath, str):
            el = get_element_by_xpath_js(driver, element_or_xpath)
            if not el:
                return False
            driver.execute_script("arguments[0].scrollIntoView(true);arguments[0].click();", el)
            return True
        else:
            driver.execute_script("arguments[0].scrollIntoView(true);arguments[0].click();", element_or_xpath)
            return True
    except Exception:
        return False


def find_in_shadow(driver, selectors):
    """在 shadow DOM 中查找元素的通用脚本：
    selectors: 列表，例如 ['#root', 'my-component', '.inner', 'div.target']。
    返回第一个匹配的节点或 None。
    """
    js = """
    const sel = arguments[0];
    function findRec(sel){
      var node = document;
      for(var i=0;i<sel.length;i++){
        var s = sel[i];
        // try querySelector on node
        try{
          var found = node.querySelector(s);
        }catch(e){var found = null}
        if(!found){
          return null;
        }
        if(i < sel.length-1){
          // go into shadowRoot if present
          node = found.shadowRoot ? found.shadowRoot : found;
        }else{
          return found;
        }
      }
      return null;
    }
    return findRec(sel);
    """
    try:
        return driver.execute_script(js, selectors)
    except Exception:
        return None


def safe_switch_to_latest_window(driver, timeout=10):
    """切换到最新打开的 window，并等待 readyState 完成与 URL 非 about:blank。
    返回 True/False 并打印简单日志（实际调用端可替换为 self.log）。"""
    try:
        handles = driver.window_handles
        if not handles:
            return False
        driver.switch_to.window(handles[-1])
        # 等待 URL 不是 about:blank 并且 document.readyState == 'complete'
        end = time.time() + timeout
        while time.time() < end:
            try:
                url = driver.current_url
                ready = driver.execute_script("return document.readyState")
            except Exception:
                url = None
                ready = None
            if url and url != 'about:blank' and ready == 'complete':
                return True
            time.sleep(0.5)
        return False
    except Exception:
        return False


def list_iframes_ids(driver):
    """返回页面上所有 iframe 的 id 列表（空 id 用索引替代）。"""
    iframes = driver.find_elements(By.TAG_NAME, 'iframe')
    res = []
    for i, f in enumerate(iframes):
        idv = f.get_attribute('id')
        res.append(idv if idv else f'index_{i}')
    return res


# ===== 示例使用（把下面几行按照注释插入你的 crawl_ppt 中适当位置） =====
# 1) 在切换到新窗口后：
# ok = safe_switch_to_latest_window(self.driver, timeout=20)
# print('switch ok', ok)
# print('current_url=', self.driver.current_url)
# print('readyState=', self.driver.execute_script("return document.readyState"))
# print('window_handles=', self.driver.window_handles)
# print('iframes=', list_iframes_ids(self.driver))
# save_debug(self.driver, prefix='after_switch')

# 2) 检查 XPath 是否存在（JS 方式）：
# xpath = '/html/body/div[1]/div[2]/div/div/div[2]/div[2]/div[3]/div[1]/div[1]/div'
# print('xpath_exists_js=', xpath_exists_js(self.driver, xpath))

# 3) 直接用 JS 取元素并点击（绕过 Selenium 的可见性/可交互限制）：
# el = get_element_by_xpath_js(self.driver, xpath)
# if el:
#     driver.execute_script("arguments[0].scrollIntoView(true);arguments[0].click();", el)
# else:
#     print('element not found by JS')

# 4) 如果元素在 iframe 中，先切换到对应 iframe：
# iframes = self.driver.find_elements(By.TAG_NAME, 'iframe')
# # 根据索引或 id 切换：
# self.driver.switch_to.frame(iframes[0])
# # 然后再调用 xpath_exists_js 或 get_element_by_xpath_js（注意：JS 查找的 document 是当前 iframe 的 document）

# 5) 如果怀疑是 shadow DOM：
# el = find_in_shadow(self.driver, ['#root', 'my-component', '.inner-class'])
# if el:
#     driver.execute_script("arguments[0].click();", el)

# 保存 debug 后，把生成的 html/png 上传我看。


def debug_find_change_item(driver, xpath=None, timeout=20, prefix="change_item_debug"):
    """在当前窗口/iframe上下文中尝试定位并调试找不到的 change-item 元素。
    返回定位到的 Selenium WebElement（在对应 iframe 中），或 None。

    使用方法：在你的 crawl_ppt 中把原来的 WebDriverWait(...) 替换为：
        el = debug_find_change_item(self.driver, xpath=你的_xpath)
        if not el:
            # 处理未找到情况
        else:
            el.click()  # 或 driver.execute_script("arguments[0].click();", el)
    """
    import time
    from selenium.webdriver.common.by import By

    if xpath is None:
        xpath = '/html/body/div[1]/div[2]/div/div/div[2]/div[2]/div[3]/div[1]/div[1]/div'

    print('=== debug_find_change_item start ===')
    try:
        print('window_handles =', driver.window_handles)
    except Exception as e:
        print('get window_handles failed:', e)

    try:
        url = driver.current_url
        print('current_url =', url)
    except Exception as e:
        print('current_url failed:', e)

    try:
        ready = driver.execute_script("return document.readyState")
        print('readyState =', ready)
    except Exception as e:
        print('readyState check failed:', e)

    # 保存调试文件
    try:
        png, html = save_debug(driver, prefix=prefix)
        print('saved debug files:', png, html)
    except Exception as e:
        print('save_debug failed:', e)

    # 列出 iframe id
    try:
        from selenium.webdriver.common.by import By as _By
        iframes = driver.find_elements(_By.TAG_NAME, 'iframe')
        ids = [f.get_attribute('id') or f.get_attribute('name') or f.get_attribute('src') or f'index_{i}' for i,f in enumerate(iframes)]
        print('iframes =', ids)
    except Exception as e:
        print('list iframes failed:', e)
        iframes = []

    # 1) 先在当前 document 用 JS 检查 XPath
    try:
        exists_now = xpath_exists_js(driver, xpath)
        print('xpath_exists_in_current_document =', exists_now)
    except Exception as e:
        print('xpath_exists_js failed:', e)
        exists_now = False

    # 2) 如果在当前 document 找到，用 JS 获取节点并返回对应的 WebElement
    if exists_now:
        try:
            el = get_element_by_xpath_js(driver, xpath)
            if el:
                print('element found by JS in current document')
                # wrap in Selenium WebElement if possible by locating via find_element (works if XPath is valid in this context)
                try:
                    we = driver.find_element(By.XPATH, xpath)
                    print('converted to Selenium WebElement via find_element')
                    return we
                except Exception:
                    # 如果转换失败，直接返回 JS 元素引用（Selenium 可以接受 JS 返回的元素用于 execute_script）
                    return el
        except Exception as e:
            print('get_element_by_xpath_js failed:', e)

    # 3) 如果当前 document 找不到，则尝试在每个 iframe 内查找
    for i, f in enumerate(iframes):
        try:
            print(f'-- try iframe index {i} --')
            driver.switch_to.default_content()
            time.sleep(0.2)
            driver.switch_to.frame(f)
            # 等待 readyState
            try:
                WebDriverWait(driver, 5).until(lambda d: d.execute_script("return document.readyState") == 'complete')
            except Exception:
                pass
            cur_url = None
            try:
                cur_url = driver.current_url
            except Exception:
                pass
            print(' iframe', i, 'url=', cur_url)

            try:
                ex = xpath_exists_js(driver, xpath)
                print(' xpath_exists_in_iframe =', ex)
            except Exception as e:
                print(' xpath_exists_js in iframe failed:', e)
                ex = False

            if ex:
                print(' element exists in iframe', i)
                try:
                    el = get_element_by_xpath_js(driver, xpath)
                    if el:
                        print('found element in iframe by JS, returning Selenium element if possible')
                        try:
                            we = driver.find_element(By.XPATH, xpath)
                            return we
                        except Exception:
                            return el
                except Exception as e:
                    print('get_element_by_xpath_js in iframe failed:', e)
        except Exception as e:
            print('error when switching to iframe', i, e)
        finally:
            try:
                driver.switch_to.default_content()
            except Exception:
                pass

    # 4) 如果仍然找不到，尝试通过 JS 直接 evaluate 并截图当前 DOM 片段
    try:
        script = (
            "var r=document.evaluate(arguments[0], document, null, XPathResult.ORDERED_NODE_SNAPSHOT_TYPE, null);"
            "var out=[]; for(var i=0;i<r.snapshotLength;i++){var n=r.snapshotItem(i); out.push(n.outerHTML ? n.outerHTML.substring(0,1000): n.tagName);} return out;"
        )
        res = driver.execute_script(script, xpath)
        print('evaluate snapshot result (outerHTML snippets):', res)
    except Exception as e:
        print('evaluate snapshot failed:', e)

    print('=== debug_find_change_item end ===')
    return None

# 说明：把原来那一段 WebDriverWait(...) 替换为：
# el = debug_find_change_item(self.driver, xpath='/html/.../div')
# if el is None:
#     # 未找到，检查生成的 debug_xxx.html/png
# else:
#     driver.execute_script('arguments[0].click();', el)






class PptCrawlerApp:
    def __init__(self, root):
        self.firstppt=True
        self.root = root
        self.root.title("课程爬取爬取")
        self.root.geometry("700x600")

        self.DRIVER_PATH = r"D:/edgedriver_win64/msedgedriver.exe"

        # Cookie 输入区
        frm_cookie = tk.Frame(self.root)
        frm_cookie.pack(pady=10, fill='x', padx=10)
        tk.Label(frm_cookie, text="请输入 Cookie：").pack(side='left')
        self.cookie_entry = tk.Entry(frm_cookie)
        self.cookie_entry.pack(side='left', fill='x', expand=True, padx=5)
        self.start_button = tk.Button(frm_cookie, text="开始获取课程列表", command=self.on_start_clicked)
        self.start_button.pack(side='left', padx=5)

        # 课程列表区
        lbl_courses = tk.Label(self.root, text="课程列表：")
        lbl_courses.pack(anchor='w', padx=10)
        self.course_listbox = tk.Listbox(self.root, height=8)
        self.course_listbox.pack(fill='x', padx=10)
        self.select_course_button = tk.Button(self.root, text="选择所选课程", command=self.on_select_course, state='disabled')
        self.select_course_button.pack(pady=5)

        # 回放列表区
        lbl_replays = tk.Label(self.root, text="回放列表（仅“回放”状态）：")
        lbl_replays.pack(anchor='w', padx=10)
        self.replay_listbox = tk.Listbox(self.root, height=8)
        self.replay_listbox.pack(fill='x', padx=10)
        self.select_replay_button = tk.Button(self.root, text="选择所选回放并开始爬取", command=self.on_select_replay, state='disabled')
        self.select_replay_button.pack(pady=5)

        #选择全部回放并开始爬取
        self.select_all_replays_button = tk.Button(self.root, text="选择全部回放并开始爬取", command=self.on_select_all_replays,
                                                   state='disabled')
        self.select_all_replays_button.pack(pady=5)

        # 进度区
        frm_progress = tk.Frame(self.root)
        frm_progress.pack(fill='x', pady=10, padx=10)
        tk.Label(frm_progress, text="爬取进度：").pack(side='left')
        self.progress_var = tk.DoubleVar(value=0.0)
        self.progress_bar = ttk.Progressbar(frm_progress, variable=self.progress_var, maximum=100)
        self.progress_bar.pack(side='left', fill='x', expand=True, padx=5)
        self.progress_label = tk.Label(frm_progress, text="0%")
        self.progress_label.pack(side='left')

        # 状态日志区
        self.status_text = tk.Text(self.root, height=8, state='disabled')
        self.status_text.pack(fill='both', padx=10, pady=5, expand=True)

        # Selenium driver placeholder
        self.driver = None
        # 存储临时用户数据目录，以便结束时清理
        self.temp_user_data_dir = None
        # 存储课程及回放信息
        self.courses = []
        self.replays = []

        # 目标 URL
        self.target_url = "https://classroom.guet.edu.cn/education/?tenant_code=21#/home?menu_code=C-wdkc,-1"
        self.log("""警告：
此应用仅可用于个人的课程学习，请尊重版权，勿将爬取内容在网上传播!""")

    def log(self, msg: str):
        """在状态框中写日志，并刷新UI"""
        self.status_text.configure(state='normal')
        self.status_text.insert('end', msg + '\n')
        self.status_text.see('end')
        self.status_text.configure(state='disabled')
        self.root.update()

    def on_start_clicked(self):
        cookie = self.cookie_entry.get().strip()
        if not cookie:
            messagebox.showerror("错误", "Cookie 不能为空！")
            return
        # 禁用按钮，防止重复点击
        self.start_button.config(state='disabled')
        self.log("开始初始化浏览器并登录...")
        self.start_crawling(cookie)

    def start_crawling(self, cookie_str: str):
        try:
            if self.firstppt is True:
                # 创建临时用户数据目录，避免与已有浏览器进程冲突
                #self.temp_user_data_dir = tempfile.mkdtemp(prefix="selenium_edge_")
                #print('开始初始化')
                edge_options = Options()
                edge_options.add_argument("--disable-blink-features=AutomationControlled")  # 避免被检测为自动化脚本
                edge_options.add_argument(
                    "user-agent=Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/132.0.0.0 Safari/537.36 Edg/132.0.0.0")  # 模拟正常浏览器的User-Agent
                edge_options.add_argument("--headless=new")  # 启用无头模式
                edge_options.add_argument("--disable-gpu")  # 禁用 GPU 加速
                edge_options.add_argument("--remote-debugging-port=9222")  # 启用远程调试端口
                edge_options.add_argument("--window-size=1920,1080")
                service = Service(self.DRIVER_PATH)
                self.driver = webdriver.Edge(service=service, options=edge_options)

                # 先打开域名主页，方便设置 cookie
                self.driver.get(self.target_url)
                time.sleep(3)
                # 解析并添加 cookie
                # 构建 cookie 列表
                #print("开始添加cookie")
                cookies = []
                for item in cookie_str.split(";"):
                    item = item.strip()  # 去掉多余的空格
                    item = item.strip('\n')
                    if "=" in item:  # 确保 item 中有 "="
                        key, value = item.split("=", 1)  # 使用 1 来限制分割成两个部分
                        cookies.append({"name": key, "value": value, "domain": "guet.edu.cn", "path": "/"})

                # 输出 cookies 列表
                #print(cookies)

                for cookie in cookies:
                    self.driver.add_cookie(cookie)


            self.driver.refresh()
            time.sleep(3)
            self.driver.get(self.target_url)
            time.sleep(3)
            #print("加载完毕")
            time.sleep(2)
            self.log("已打开目标页面，切换 iframe ...")
            # 切换到 iframe
            try:
                iframe = WebDriverWait(self.driver, 10).until(
                    EC.presence_of_element_located((By.ID, "inlineFrameExample"))
                )
                self.driver.switch_to.frame(iframe)
            except Exception as e:
                self.log(f"未能切换到 iframe: {e}")
                self.cleanup_driver()
                self.enable_start_button()
                return

            # 等待并获取课程元素
            elems = WebDriverWait(self.driver, 10).until(
                EC.presence_of_all_elements_located((By.CSS_SELECTOR, "div.courseInfo"))
            )
            self.courses.clear()
            for idx, ce in enumerate(elems):
                try:
                    title = ce.find_element(By.CSS_SELECTOR, ".course-label").text.strip()
                    teacher = ce.find_element(By.CSS_SELECTOR, ".course-desc-p").text.strip()
                    self.courses.append({"title": title, "teacher": teacher, "element": ce})
                    self.log(f"课程 {idx+1}: {title} - {teacher}")
                except Exception:
                    continue

            # 更新 GUI 课程列表
            self.course_listbox.delete(0, 'end')
            for c in self.courses:
                self.course_listbox.insert('end', f"{c['title']} - {c['teacher']}")
            self.select_course_button.config(state='normal')
            self.log("请在列表中选择课程，然后点击“选择所选课程”")
        except Exception as ex:
            self.log(f"初始化或获取课程列表时出错: {ex}")
            self.cleanup_driver()
            self.enable_start_button()

    def enable_start_button(self):
        self.start_button.config(state='normal')
        self.root.update()

    def on_select_course(self):
        sel = self.course_listbox.curselection()
        if not sel:
            messagebox.showerror("错误", "请先选择课程！")
            return
        index = sel[0]
        selected = self.courses[index]
        self.log(f"选择了课程：{selected['title']} - {selected['teacher']}")
        self.select_course_button.config(state='disabled')
        self.fetch_replays(selected)

    def fetch_replays(self, course_item: dict):
        time.sleep(3)
        try:
            # 点击进入课程
            try:
                btn = course_item["element"].find_element(By.CSS_SELECTOR, ".el-button--primary")
                btn.click()
            except Exception as e:
                self.log(f"点击课程进入失败: {e}")
                self.select_course_button.config(state='normal')
                return
            time.sleep(2)
            # 切换到新窗口
            handles = self.driver.window_handles
            if len(handles) > 1:
                self.driver.switch_to.window(handles[-1])
            time.sleep(1)

            # 获取所有 <p> 标签，即回放记录
            p_elements = self.driver.find_elements(By.CSS_SELECTOR, ".content-inner-one > p")
            self.replays.clear()
            for idx, p in enumerate(p_elements):
                try:
                    spans = p.find_elements(By.TAG_NAME, "span")
                    time_info = spans[1].text.strip()
                    teacher = spans[2].get_attribute("title") or spans[2].text.strip()
                    status = spans[4].text.strip()
                    self.replays.append({"time": time_info, "teacher": teacher, "status": status, "element": p})
                    self.log(f"回放 {idx+1}: {time_info} - {teacher} ({status})")
                except Exception:
                    continue

            # 更新 GUI 回放列表，仅显示 status == "回放"
            self.replay_listbox.delete(0, 'end')
            for r in self.replays:
                if r["status"] == "回放":
                    self.replay_listbox.insert('end', f"{r['time']} - {r['teacher']}")
            self.select_replay_button.config(state='normal')
            self.select_all_replays_button.config(state='normal')
            self.log("请在列表中选择回放，然后点击“选择所选回放并开始爬取”")
        except Exception as ex:
            self.log(f"获取回放列表出错: {ex}")
            self.select_course_button.config(state='normal')

    def on_select_replay(self):
        sel = self.replay_listbox.curselection()
        if not sel:
            messagebox.showerror("错误", "请先选择回放！")
            return
        index = sel[0]
        # 过滤后的列表
        filtered = [r for r in self.replays if r["status"] == "回放"]
        selected = filtered[index]
        self.log(f"选择了回放：{selected['time']} - {selected['teacher']}")
        self.select_replay_button.config(state='disabled')
        self.crawl_ppt(selected)

    def crawl_ppt(self, replay_item: dict):

        # 点击回放按钮打开回放窗口
        try:
            parent = self.driver.current_window_handle
            #print('开始回放按钮')
            spans = replay_item["element"].find_elements(By.TAG_NAME, "span")
            btn = spans[4]
            #print('回放按钮查找成功，进行点击')
            btn.click()
            time.sleep(3)
        except Exception as e:
            self.log(f"点击回放进入失败: {e}")
            self.select_replay_button.config(state='normal')
            return
        # 3. 切到新窗口（总是最后一个）
        # WebDriverWait(self.driver, 10).until(lambda d: len(d.window_handles) > 1)
        # new_win = [h for h in self.driver.window_handles if h != parent][0]
        # self.driver.switch_to.window(new_win)

        self.driver.switch_to.window(self.driver.window_handles[-1])

        time.sleep(3)

        # 4. 在新窗口做完事以后，先关闭它
        self.driver.close()

        # 5. 切回父窗口
        #self.driver.switch_to.window(parent)
        self.driver.switch_to.window(self.driver.window_handles[-1])


        # 6. 重新定位按钮并再次点击
        status_button = replay_item["element"].find_elements(By.TAG_NAME, "span")[4]
        #print("准备再次点击")
        status_button.click()
        #print("准备再次点击成功")
        time.sleep(10)
        # 切换到新打开的窗口/标签
        all_windows = self.driver.window_handles  # 获取所有窗口句柄
        self.driver.switch_to.window(all_windows[-1])  # 切换到最新打开的窗口
        time.sleep(3)
        # 准备 PPT
        prs = Presentation()

        # 1) 在切换到新窗口后：
        ok = safe_switch_to_latest_window(self.driver, timeout=20)
        print('switch ok', ok)
        print('current_url=', self.driver.current_url)
        print('readyState=', self.driver.execute_script("return document.readyState"))
        print('window_handles=', self.driver.window_handles)
        print('iframes=', list_iframes_ids(self.driver))
        #save_debug(self.driver, prefix='after_switch')

        # 2) 检查 XPath 是否存在（JS 方式）：
        xpath = '//div[@class="change-item__img"]'
        print('xpath_exists_js=', xpath_exists_js(self.driver, xpath))

        el = get_element_by_xpath_js(self.driver, xpath)
        self.driver.execute_script("arguments[0].scrollIntoView(true);arguments[0].click();", el)

        # 获取总页数信息
        try:
            page_info_elem = WebDriverWait(self.driver, 10).until(
                EC.presence_of_element_located((By.CLASS_NAME, 'ppt_page_con'))
            )
            page_info = page_info_elem.text.strip()
            current_page, total_pages = map(int, page_info.split('/'))
        except Exception:
            self.log("无法获取页数信息，停止爬取")
            self.cleanup_driver()
            self.enable_start_button()
            return

        # 从第一页开始循环抓取
        for page_no in range(1, total_pages + 1):
            time.sleep(2)
            # 等待 canvas 出现并抓图
            try:
                WebDriverWait(self.driver, 10).until(
                    EC.presence_of_element_located((By.ID, 'ppt_canvas'))
                )
                img_data = self.driver.execute_script(
                    "return document.getElementById('ppt_canvas').toDataURL('image/png').substring(22);"
                )
                img_bytes = base64.b64decode(img_data)
                slide_layout = prs.slide_layouts[5]
                slide = prs.slides.add_slide(slide_layout)
                image_stream = BytesIO(img_bytes)
                # left = Inches(1); top = Inches(1); height = Inches(5)
                # slide.shapes.add_picture(image_stream, left, top, height=height)
                pic = slide.shapes.add_picture(image_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)
                self.log(f"已抓取第 {page_no} 页 / 共 {total_pages} 页")
            except Exception as e:
                self.log(f"第 {page_no} 页抓取失败: {e}")
            # 更新进度条
            pct = page_no / total_pages * 100
            self.progress_var.set(pct)
            self.progress_label.config(text=f"{int(pct)}%")
            self.root.update()
            # 如果不是最后一页，点击“下一页”
            if page_no < total_pages:
                try:
                    next_btn = self.driver.find_element(By.CLASS_NAME, 'ppt_btn_next')
                    next_btn.click()
                    time.sleep(0.1)
                except Exception as e:
                    self.log(f"翻页到第 {page_no+1} 页失败: {e}")
                    break

        # 提示用户输入文件名
        filename = simpledialog.askstring("保存文件", "请输入 PPT 文件名（不带后缀）：", parent=self.root)
        if not filename:
            filename = f"output_{int(time.time())}"
        save_path = f"PPT\\{filename}.pptx"
        prs.save(save_path)
        self.log(f"PPT 已保存为: {save_path}")
        messagebox.showinfo("完成", f"PPT 已保存为当前目录下的：{save_path}")

        # 结束时关闭浏览器并清理
        # self.cleanup_driver()

        # 重置 UI 状态
        self.progress_var.set(0)
        self.progress_label.config(text="0%")
        self.start_button.config(state='active')
        self.select_course_button.config(state='disabled')
        self.select_replay_button.config(state='disabled')
        self.course_listbox.delete(0, 'end')
        self.replay_listbox.delete(0, 'end')
        self.firstppt=False
        self.log("流程结束，可再次重试")



        # self.driver.quit()
        # os.system("taskkill /F /IM msedge.exe")
        # # 然后结束主循环并退出程序
        # self.root.destroy()
        # # 如果希望彻底退出，避免残留线程，可再调用 sys.exit()
        # sys.exit(0)

        self.driver.close()
        self.driver.switch_to.window(self.driver.window_handles[-1])
        self.driver.close()
        self.driver.switch_to.window(self.driver.window_handles[-1])

    def on_select_all_replays(self):
        """选择全部回放并开始爬取"""
        if not self.replays:
            messagebox.showerror("错误", "没有回放可供选择！")
            return

        # 禁用按钮，防止重复点击
        self.select_replay_button.config(state='disabled')
        self.select_course_button.config(state='disabled')

        # 循环爬取所有回放
        print("开始循环爬取所有回放")
        self.replay_len=0
        self.replay_idx=0
        for replay_item in self.replays:
            if replay_item["status"] == "回放":  # 只爬取状态为回放的项
                self.replay_len=self.replay_len+1
        for replay_item in self.replays:
            if replay_item["status"] == "回放":  # 只爬取状态为回放的项
                print(f"开始爬取回放：{replay_item['time']} - {replay_item['teacher']}")
                self.log(f"开始爬取回放：{replay_item['time']} - {replay_item['teacher']}")
                self.crawl_ppt_for_all(replay_item)
                self.replay_idx=self.replay_idx+1


        # 提示用户爬取完成
        messagebox.showinfo("完成", "所有回放已爬取完成！")
        self.log("所有回放爬取完成。")

        # 重置 UI 状态
        self.progress_var.set(0)
        self.progress_label.config(text="0%")
        self.select_replay_button.config(state='normal')
        self.select_course_button.config(state='normal')
        self.course_listbox.delete(0, 'end')
        self.replay_listbox.delete(0, 'end')

    def crawl_ppt_for_all(self, replay_item: dict):
        """爬取单个回放并自动生成文件名（完整、鲁棒版）。
        替换原有函数即可：使用 locate_replay_parent 和 safe_click_span_button，每次重新定位元素，处理窗口切换与异常。
        """
        try:
            filename = f"{replay_item['time']}_{replay_item['teacher']}.pptx"
            save_path = f"PPT\\{filename}"
            print(save_path)

            # ---------- 1) 找到或确定“回放列表所在窗口” ----------
            # 优先使用已有的 self.replay_window（fetch_replays 在第一次抓取时会设置）
            replay_window = getattr(self, 'replay_window', None)
            if replay_window not in (self.driver.window_handles if self.driver else []):
                # 如果没有或失效，尝试在所有窗口中找到包含回放列表的窗口
                found = False
                for h in list(self.driver.window_handles):
                    try:
                        self.driver.switch_to.window(h)
                        try:
                            self.driver.switch_to.default_content()
                        except Exception:
                            pass
                        # 检查页面是否有回放列表的标识（p 元素）
                        elems = self.driver.find_elements(By.CSS_SELECTOR, ".content-inner-one > p")
                        if elems and len(elems) > 0:
                            replay_window = h
                            self.replay_window = h
                            found = True
                            break
                    except Exception:
                        continue
                if not found:
                    # 仍未找到，使用当前窗口作为回退
                    replay_window = self.driver.current_window_handle
                    self.replay_window = replay_window

            # 切回回放列表窗口并确保 default_content
            try:
                self.driver.switch_to.window(replay_window)
                self.driver.switch_to.default_content()
            except Exception:
                # 尝试切回最后一个窗口
                self.driver.switch_to.window(self.driver.window_handles[-1])
                try:
                    self.driver.switch_to.default_content()
                except Exception:
                    pass

            # 等待回放列表元素至少存在（短等待）
            try:
                WebDriverWait(self.driver, 6).until(
                    EC.presence_of_element_located((By.CSS_SELECTOR, ".content-inner-one > p"))
                )
            except Exception:
                # 继续也行，后面 locate_replay_parent 会再尝试
                pass

            # ---------- 2) 初次定位该回放条目并点击（打开临时窗口并关闭） ----------
            parent = locate_replay_parent(self, replay_item['time'], replay_item['teacher'], timeout=6)
            if parent is None:
                self.log(f"无法重新定位回放项：{replay_item['time']} - {replay_item['teacher']}")
                return

            # 记录当前窗口集合数量
            prev_handles = set(self.driver.window_handles)
            ok = safe_click_span_button(self, parent, span_index=4, timeout=6)
            if not ok:
                self.log(f"点击回放进入失败（首次点击）：{replay_item['time']} - {replay_item['teacher']}")
                return

            # 等待新窗口出现（短超时），若出现则关闭该临时窗口并回到回放列表
            try:
                WebDriverWait(self.driver, 8).until(lambda d: len(d.window_handles) > len(prev_handles))
                # 切换到最新打开的窗口（临时弹窗）
                newest = [h for h in self.driver.window_handles if h not in prev_handles]
                if newest:
                    newh = newest[-1]
                    try:
                        self.driver.switch_to.window(newh)
                        time.sleep(1)
                        # 直接关闭它（这是你原流程的第一步）
                        try:
                            self.driver.close()
                        except Exception:
                            pass
                    except Exception:
                        pass
                # 切回回放列表窗口
                try:
                    self.driver.switch_to.window(self.replay_window)
                except Exception:
                    if self.driver.window_handles:
                        self.driver.switch_to.window(self.driver.window_handles[-1])
            except Exception:
                # 没有新窗口也继续（可能第一次点击就直接进入播放窗口或没有新弹窗）
                try:
                    self.driver.switch_to.window(self.replay_window)
                except Exception:
                    pass

            time.sleep(0.8)

            # ---------- 3) 再次定位并点击（正式打开播放窗口） ----------
            parent = locate_replay_parent(self, replay_item['time'], replay_item['teacher'], timeout=6)
            if parent is None:
                self.log(f"二次定位回放项失败（已尝试重定位）：{replay_item['time']} - {replay_item['teacher']}")
                return

            prev_handles = set(self.driver.window_handles)
            ok = safe_click_span_button(self, parent, span_index=4, timeout=8)
            if not ok:
                # 尝试用 JS xpath 最后手段点击
                try:
                    js_xpath = f"//p[span[2][contains(text(), '{replay_item['time']}')] and span[3][contains(text(), '{replay_item['teacher']}')]]//span[5]"
                    el = get_element_by_xpath_js(self.driver, js_xpath)
                    if el:
                        self.driver.execute_script("arguments[0].click();", el)
                        ok = True
                except Exception:
                    ok = False

            if not ok:
                self.log("尝试多种方法点击回放按钮均失败，跳过该回放。")
                return

            # 等待播放窗口真正出现并切换到它
            try:
                WebDriverWait(self.driver, 12).until(lambda d: len(d.window_handles) > len(prev_handles))
                # 切换到最新打开的窗口
                new_handles = [h for h in self.driver.window_handles if h not in prev_handles]
                if not new_handles:
                    # 兜底：切到最后一个
                    self.driver.switch_to.window(self.driver.window_handles[-1])
                else:
                    self.driver.switch_to.window(new_handles[-1])
            except Exception:
                # 如果超时，也尝试切换到最后一个窗口
                try:
                    self.driver.switch_to.window(self.driver.window_handles[-1])
                except Exception:
                    pass

            time.sleep(2)

            # ---------- 4) 在播放窗口中准备抓 PPT ----------
            prs = Presentation()
            try:
                ok = safe_switch_to_latest_window(self.driver, timeout=20)
            except Exception:
                ok = False
            print('switch ok', ok)
            try:
                print('current_url=', self.driver.current_url)
            except Exception:
                pass
            try:
                print('readyState=', self.driver.execute_script("return document.readyState"))
            except Exception:
                pass
            print('window_handles=', self.driver.window_handles)
            print('iframes=', list_iframes_ids(self.driver))

            # 点击进入 PPT 视图（change-item__img）
            xpath = '//div[@class="change-item__img"]'
            try:
                # 先用 JS 检查是否存在，再点击（更稳）
                if xpath_exists_js(self.driver, xpath):
                    el = get_element_by_xpath_js(self.driver, xpath)
                    if el:
                        try:
                            self.driver.execute_script("arguments[0].scrollIntoView(true);arguments[0].click();", el)
                        except Exception:
                            # 兜底：用 Selenium find + click
                            try:
                                we = self.driver.find_element(By.XPATH, xpath)
                                we.click()
                            except Exception:
                                pass
                    else:
                        # 兜底：Selenium 查找并点击
                        try:
                            we = WebDriverWait(self.driver, 6).until(EC.element_to_be_clickable((By.XPATH, xpath)))
                            we.click()
                        except Exception:
                            pass
                else:
                    # 如果没有 change-item，继续（有些页面直接显示 canvas）
                    pass
            except Exception:
                pass

            # 等待页数信息
            try:
                page_info_elem = WebDriverWait(self.driver, 12).until(
                    EC.presence_of_element_located((By.CLASS_NAME, 'ppt_page_con'))
                )
                page_info = page_info_elem.text.strip()
                current_page, total_pages = map(int, page_info.split('/'))
            except Exception:
                self.log("无法获取页数信息，停止爬取该条回放")
                # 在出错时关闭播放窗口并返回回放列表窗口
                try:
                    self.driver.close()
                except Exception:
                    pass
                try:
                    if self.replay_window in self.driver.window_handles:
                        self.driver.switch_to.window(self.replay_window)
                    elif self.driver.window_handles:
                        self.driver.switch_to.window(self.driver.window_handles[-1])
                except Exception:
                    pass
                return

            # ---------- 5) 循环抓取每一页 canvas 并保存到 ppt ----------
            for page_no in range(1, total_pages + 1):
                try:
                    # 等待 canvas 出现
                    WebDriverWait(self.driver, 12).until(
                        EC.presence_of_element_located((By.ID, 'ppt_canvas'))
                    )
                    # 取图片 base64（若失败重试一次）
                    try:
                        img_data = self.driver.execute_script(
                            "return document.getElementById('ppt_canvas').toDataURL('image/png').substring(22);"
                        )
                    except Exception:
                        time.sleep(0.5)
                        img_data = self.driver.execute_script(
                            "return document.getElementById('ppt_canvas').toDataURL('image/png').substring(22);"
                        )
                    img_bytes = base64.b64decode(img_data)
                    slide_layout = prs.slide_layouts[5]
                    slide = prs.slides.add_slide(slide_layout)
                    image_stream = BytesIO(img_bytes)
                    pic = slide.shapes.add_picture(image_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)
                    self.log(f"已抓取第 {page_no} 页 / 共 {total_pages} 页")
                except Exception as e:
                    self.log(f"第 {page_no} 页抓取失败: {e}")

                # 更新进度条（按你原来的复合比例）
                try:
                    pct = (self.replay_idx / self.replay_len + page_no / total_pages / self.replay_len) * 100
                except Exception:
                    pct = page_no / total_pages * 100
                self.progress_var.set(pct)
                self.progress_label.config(text=f"{int(pct)}%")
                self.root.update()

                # 如果不是最后一页，点击“下一页”
                if page_no < total_pages:
                    clicked = False
                    try:
                        next_btn = WebDriverWait(self.driver, 6).until(
                            EC.element_to_be_clickable((By.CLASS_NAME, 'ppt_btn_next'))
                        )
                        try:
                            next_btn.click()
                            clicked = True
                        except Exception:
                            # JS 点击兜底
                            try:
                                self.driver.execute_script("arguments[0].scrollIntoView(true);arguments[0].click();",
                                                           next_btn)
                                clicked = True
                            except Exception:
                                clicked = False
                    except Exception:
                        # 仍可尝试直接用 JS 点击 known xpath/class
                        try:
                            self.driver.execute_script(
                                "var b=document.getElementsByClassName('ppt_btn_next')[0]; if(b){b.scrollIntoView();b.click();}"
                            )
                            clicked = True
                        except Exception:
                            clicked = False

                    if not clicked:
                        self.log(f"翻页到第 {page_no + 1} 页失败: 无法点击下一页按钮，停止本回放抓取")
                        break
                    # 给一点渲染时间
                    time.sleep(4)

            # ---------- 6) 保存 PPT 并清理窗口 ----------
            try:
                # 确保目录存在
                os.makedirs(os.path.dirname(save_path), exist_ok=True)
                prs.save(save_path)
                self.log(f"PPT 已保存为: {save_path}")
            except Exception as e:
                self.log(f"PPT 保存失败: {e}")

            # 关闭播放窗口并切回回放列表窗口
            try:
                self.driver.close()
            except Exception:
                pass
            try:
                if hasattr(self, 'replay_window') and self.replay_window in self.driver.window_handles:
                    self.driver.switch_to.window(self.replay_window)
                elif self.driver.window_handles:
                    self.driver.switch_to.window(self.driver.window_handles[-1])
            except Exception:
                pass

            time.sleep(1)

        except Exception as e:
            # 任意未捕获异常时记录并尝试恢复到回放列表窗口
            self.log(f"crawl_ppt_for_all 失败: {e}")
            try:
                if hasattr(self, 'replay_window') and self.replay_window in self.driver.window_handles:
                    self.driver.switch_to.window(self.replay_window)
            except Exception:
                pass
            return


    def cleanup_driver(self):
        # 关闭浏览器并删除临时用户数据目录
        try:
            if self.driver:
                self.driver.quit()
                os.system("taskkill /F /IM msedge.exe")
                # 然后结束主循环并退出程序
                self.root.destroy()
                # 如果希望彻底退出，避免残留线程，可再调用 sys.exit()
                sys.exit(0)
        except:
            pass

def main():
    root = tk.Tk()
    PptCrawlerApp(root)  # 不必保存到变量也能正常工作
    root.mainloop()

if __name__ == "__main__":
    main()
